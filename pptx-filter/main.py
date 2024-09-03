from pptx import Presentation

PRESENTATION_PATH = 'presentations/'
COPY_PRESENTATION_PATH = 'copy/'
PPTX_FILE = 'your_presentation.pptx'
COPY_FILE = 'copy.pptx'
FILTER_KEYWORDS = ['keyword1', 'keyword2', 'keyword3']

""" check if the extracted text is valid """
def extracted_text_linter(notes_text: str) -> str:

    temp_text = notes_text.split(', ')

    for text in temp_text:

        if not text.isalpha():

            return False
        
    return True

""" Extract and hide slides based on keywords """
def extract_and_hide_slides(pptx_file: str, new_presentation_name: str, filter_keywords: list) -> None:
    
    # Open the presentation
    presentation = Presentation(PRESENTATION_PATH + pptx_file)
    
    # Copy this presentation to a new presentation
    new_presentation = presentation

    # Keep slides that either contain the keywords or no keywords at all
    slides_to_keep = []

    # Loop through each slide
    for slide in presentation.slides:

        notes_slide = slide.notes_slide
        keywords = []
        
        if notes_slide:

            notes_text = notes_slide.notes_text_frame.text

            while 'Market=[' in notes_text:

                index = notes_text.find('[')

                if index != -1:

                    extracted_text = notes_text[index+1:notes_text.find(']')]

                    if extracted_text_linter(extracted_text):

                        pass

                    else:

                        print('Invalid text found:', extracted_text)

                    keywords += extracted_text.split(', ')
                    notes_text = notes_text[notes_text.find(']')+1:]

            keywords_found = [keyword for keyword in keywords if keyword in filter_keywords]

            if not keywords:

                slides_to_keep.append(slide)
            
            elif keywords_found:

                print('Slide with keywords found:', keywords_found)
                slides_to_keep.append(slide)

    print('Slides to keep:', slides_to_keep)

    # delete slides from new presentation that are not in slides_to_keep
    for slide in new_presentation.slides:

        if slide not in slides_to_keep:
            
            slide_index = new_presentation.slides.index(slide)
            del new_presentation.slides._sldIdLst[slide_index]

    # save the new presentation
    new_presentation.save(COPY_PRESENTATION_PATH + new_presentation_name)

extract_and_hide_slides(PPTX_FILE, COPY_FILE, FILTER_KEYWORDS)
